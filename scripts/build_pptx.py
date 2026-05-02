#!/usr/bin/env python3
"""
Assemble an editable .pptx from a slides.json plan plus a directory of SVGs.

Each slide:
  - full-bleed SVG background (image/svg+xml registered so PowerPoint keeps it
    editable; right-click -> "Convert to Shape" turns vector SVGs into native
    PPT shapes).
  - foreground title + bullet text frames.
  - footer with deck title + slide number.
  - speaker_notes copied into the notes pane.

Usage:
    python build_pptx.py \\
        --slides analysis/slides.json \\
        --svg-dir images/svg \\
        --out out/topic.pptx \\
        [--deck-title "Deck Title"] \\
        [--template template.pptx]

Deps: python-pptx, lxml.   pip install python-pptx lxml
"""
from __future__ import annotations

import argparse
import json
import sys
from pathlib import Path

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.oxml.ns import qn
    from lxml import etree
except ImportError as e:
    print(
        "Missing deps. Install with:\n  pip install python-pptx lxml\n"
        f"Original error: {e}",
        file=sys.stderr,
    )
    sys.exit(2)


# ---------- SVG content-type registration ----------
SVG_RELATIONSHIP = "http://schemas.openxmlformats.org/officeDocument/2006/relationships/image"
SVG_CT = "image/svg+xml"


def register_svg_content_type(prs: Presentation) -> None:
    """Ensure [Content_Types].xml has a Default for svg -> image/svg+xml."""
    ct_part = prs.part.package._content_types  # internal but stable in python-pptx
    # python-pptx exposes a dict-like; safest is to operate on its element.
    try:
        ct_elem = ct_part._element  # lxml Element
    except AttributeError:
        ct_elem = None

    if ct_elem is None:
        # Fallback: walk parts for [Content_Types].xml
        for part in prs.part.package.iter_parts():
            if part.partname == "/[Content_Types].xml":
                ct_elem = etree.fromstring(part.blob)
                break

    if ct_elem is None:
        return  # best effort

    ns = "http://schemas.openxmlformats.org/package/2006/content-types"
    has_svg = any(
        d.get("Extension", "").lower() == "svg"
        for d in ct_elem.findall(f"{{{ns}}}Default")
    )
    if not has_svg:
        default = etree.SubElement(ct_elem, f"{{{ns}}}Default")
        default.set("Extension", "svg")
        default.set("ContentType", SVG_CT)


# ---------- helpers ----------
def add_svg_picture(slide, svg_path: Path, left, top, width, height):
    """
    Insert an SVG as a picture. python-pptx's add_picture infers content type
    from extension; with the Default override above, .svg becomes image/svg+xml.
    PowerPoint shows the SVG; right-click -> Convert to Shape makes it editable.
    """
    return slide.shapes.add_picture(
        str(svg_path), left, top, width=width, height=height
    )


def parse_palette(palette_str: str):
    """Best-effort: return (primary, accent, neutral) RGB tuples from a free-form
    palette string. Falls back to a sober default if no hex codes are found."""
    import re

    hexes = re.findall(r"#([0-9a-fA-F]{6})", palette_str or "")
    rgbs = [tuple(int(h[i : i + 2], 16) for i in (0, 2, 4)) for h in hexes]
    while len(rgbs) < 3:
        rgbs.append([(11, 37, 69), (67, 176, 241), (245, 247, 250)][len(rgbs)])
    return rgbs[0], rgbs[1], rgbs[2]


def add_text(slide, left, top, width, height, text, *, size=18, bold=False, rgb=(20, 20, 20), align_left=True):
    tx = slide.shapes.add_textbox(left, top, width, height)
    tf = tx.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    if not align_left:
        from pptx.enum.text import PP_ALIGN
        p.alignment = PP_ALIGN.RIGHT
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = RGBColor(*rgb)
    return tx


def add_bullets(slide, left, top, width, height, bullets, *, rgb=(40, 40, 40)):
    if not bullets:
        return None
    tx = slide.shapes.add_textbox(left, top, width, height)
    tf = tx.text_frame
    tf.word_wrap = True
    for i, b in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        run = p.add_run()
        run.text = "•  " + b
        run.font.size = Pt(18)
        run.font.color.rgb = RGBColor(*rgb)
        p.space_after = Pt(6)
    return tx


def set_speaker_notes(slide, text: str) -> None:
    if not text:
        return
    notes = slide.notes_slide
    tf = notes.notes_text_frame
    tf.text = text


# ---------- main ----------
def build(
    slides_json: Path,
    svg_dir: Path,
    out_path: Path,
    deck_title: str,
    template: Path | None,
) -> None:
    plan = json.loads(slides_json.read_text(encoding="utf-8"))
    if not plan:
        raise SystemExit("slides.json is empty")

    prs = Presentation(str(template)) if template and template.exists() else Presentation()
    # 16:9
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    register_svg_content_type(prs)

    palette = plan[0].get("visual", {}).get("palette", "")
    primary, accent, neutral = parse_palette(palette)

    blank_layout = prs.slide_layouts[6] if len(prs.slide_layouts) > 6 else prs.slide_layouts[0]

    SW, SH = prs.slide_width, prs.slide_height
    margin = Inches(0.5)

    for entry in plan:
        n = entry.get("n")
        title = entry.get("title", "")
        subtitle = entry.get("subtitle", "")
        bullets = entry.get("bullets", []) or []
        notes = entry.get("speaker_notes", "")

        slide = prs.slides.add_slide(blank_layout)

        # Background: full-bleed SVG
        svg_path = svg_dir / f"slide-{int(n):02d}.svg"
        if not svg_path.exists():
            # try looser names
            candidates = list(svg_dir.glob(f"slide-{n}*.svg")) + list(
                svg_dir.glob(f"slide-{int(n):02d}*.svg")
            )
            svg_path = candidates[0] if candidates else None  # type: ignore

        if svg_path and svg_path.exists():
            try:
                add_svg_picture(slide, svg_path, 0, 0, SW, SH)
            except Exception as e:
                print(f"[warn] slide {n}: SVG insert failed ({e}); skipping background")

        # Title
        add_text(
            slide,
            margin,
            margin,
            SW - 2 * margin,
            Inches(0.9),
            title,
            size=32,
            bold=True,
            rgb=primary,
        )
        if subtitle:
            add_text(
                slide,
                margin,
                margin + Inches(0.85),
                SW - 2 * margin,
                Inches(0.5),
                subtitle,
                size=16,
                bold=False,
                rgb=accent,
            )

        # Bullets — left third by default
        if bullets:
            add_bullets(
                slide,
                margin,
                Inches(2.0),
                Inches(5.5),
                Inches(4.5),
                bullets,
                rgb=(40, 40, 40),
            )

        # Footer
        add_text(
            slide,
            margin,
            SH - Inches(0.4),
            Inches(8),
            Inches(0.3),
            deck_title,
            size=10,
            rgb=(120, 120, 120),
        )
        add_text(
            slide,
            SW - Inches(1.5),
            SH - Inches(0.4),
            Inches(1),
            Inches(0.3),
            f"{n} / {len(plan)}",
            size=10,
            rgb=(120, 120, 120),
            align_left=False,
        )

        set_speaker_notes(slide, notes)

    out_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out_path))
    print(f"wrote {out_path} ({len(plan)} slides)")


def main() -> int:
    ap = argparse.ArgumentParser(description="Build editable PPTX from slides.json + SVGs")
    ap.add_argument("--slides", type=Path, required=True)
    ap.add_argument("--svg-dir", type=Path, required=True)
    ap.add_argument("--out", type=Path, required=True)
    ap.add_argument("--deck-title", default="")
    ap.add_argument("--template", type=Path, default=None)
    args = ap.parse_args()

    deck_title = args.deck_title or args.out.stem
    build(args.slides, args.svg_dir, args.out, deck_title, args.template)
    return 0


if __name__ == "__main__":
    sys.exit(main())
